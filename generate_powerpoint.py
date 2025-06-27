import csv
import sys
from pptx import Presentation
from pptx.util import Inches

def generate_powerpoint(csv_file_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Technical SEO Audit Report"
    subtitle.text = "Generated from Screaming Frog Data"

    # Read the CSV and process issues
    issues = []
    try:
        with open(csv_file_path, mode='r', encoding='utf-8') as file:
            reader = csv.DictReader(file)
            for row in reader:
                issues.append(row)
    except FileNotFoundError:
        print(f"Error: CSV file not found at {csv_file_path}")
        sys.exit(1)
    except Exception as e:
        print(f"Error reading CSV file: {e}")
        sys.exit(1)

    if not issues:
        print("No issues found in the CSV to generate a report.")
        sys.exit(0)

    # Issues Overview Slide
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Issues Overview"

    # Group issues by category (e.g., 'Issue')
    issue_categories = {}
    for issue in issues:
        issue_name = issue.get('Issue', 'Unknown Issue')
        if issue_name not in issue_categories:
            issue_categories[issue_name] = []
        issue_categories[issue_name].append(issue)

    body_text = ""
    for category, items in issue_categories.items():
        body_text += f"\n- {category}: {len(items)} occurrences"

    body.text = body_text.strip()

    # Individual Issue Slides (example: top 5 issues)
    sorted_categories = sorted(issue_categories.items(), key=lambda item: len(item[1]), reverse=True)
    for category, items in sorted_categories[:5]: # Limit to top 5 for brevity
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = f"Issue: {category} ({len(items)} occurrences)"
        
        # Display first few URLs for context
        detail_text = "Affected URLs:\n"
        for i, item in enumerate(items[:10]): # Show up to 10 affected URLs
            detail_text += f"- {item.get('Address', 'N/A')}\n"
        if len(items) > 10:
            detail_text += f"...and {len(items) - 10} more."
        
        body.text = detail_text.strip()

    # Save the presentation
    output_pptx_path = csv_file_path.replace('.csv', '_report.pptx')
    prs.save(output_pptx_path)
    print(f"PowerPoint presentation generated successfully: {output_pptx_path}")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python3 generate_powerpoint.py <path_to_issues_overview_csv>")
        sys.exit(1)
    
    csv_file = sys.argv[1]
    generate_powerpoint(csv_file)
